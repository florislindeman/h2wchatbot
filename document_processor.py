from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import pytesseract
import io
import logging
from typing import BinaryIO

logger = logging.getLogger(__name__)

class DocumentProcessor:
    
    @staticmethod
    def extract_text_from_pdf(file_content: bytes) -> str:
        """Extract text from PDF"""
        try:
            pdf_file = io.BytesIO(file_content)
            reader = PdfReader(pdf_file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            logger.info(f"Extracted {len(text)} characters from PDF")
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_docx(file_content: bytes) -> str:
        """Extract text from DOCX"""
        try:
            docx_file = io.BytesIO(file_content)
            doc = DocxDocument(docx_file)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            logger.info(f"Extracted {len(text)} characters from DOCX")
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_xlsx(file_content: bytes) -> str:
        """Extract text from Excel"""
        try:
            xlsx_file = io.BytesIO(file_content)
            workbook = load_workbook(xlsx_file, data_only=True)
            text = ""
            
            for sheet in workbook.worksheets:
                text += f"\n=== Sheet: {sheet.title} ===\n"
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    text += row_text + "\n"
            
            logger.info(f"Extracted {len(text)} characters from XLSX")
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting XLSX text: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_pptx(file_content: bytes) -> str:
        """Extract text from PowerPoint"""
        try:
            pptx_file = io.BytesIO(file_content)
            presentation = Presentation(pptx_file)
            text = ""
            
            for i, slide in enumerate(presentation.slides):
                text += f"\n=== Slide {i+1} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            logger.info(f"Extracted {len(text)} characters from PPTX")
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_image(file_content: bytes) -> str:
        """Extract text from image using OCR"""
        try:
            image = Image.open(io.BytesIO(file_content))
            text = pytesseract.image_to_string(image, lang='nld+eng')
            logger.info(f"Extracted {len(text)} characters from image via OCR")
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting text from image: {e}")
            return ""
    
    @staticmethod
    def process_document(file_content: bytes, file_type: str) -> str:
        """Process document based on file type and extract text"""
        file_type = file_type.lower()
        
        if file_type == 'pdf':
            return DocumentProcessor.extract_text_from_pdf(file_content)
        elif file_type in ['docx', 'doc']:
            return DocumentProcessor.extract_text_from_docx(file_content)
        elif file_type in ['xlsx', 'xls']:
            return DocumentProcessor.extract_text_from_xlsx(file_content)
        elif file_type in ['pptx', 'ppt']:
            return DocumentProcessor.extract_text_from_pptx(file_content)
        elif file_type in ['jpg', 'jpeg', 'png', 'gif']:
            return DocumentProcessor.extract_text_from_image(file_content)
        else:
            logger.warning(f"Unsupported file type: {file_type}")
            return ""
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1500) -> list[str]:
        """Split text into chunks of approximately chunk_size words"""
        words = text.split()
        chunks = []
        current_chunk = []
        current_size = 0
        
        for word in words:
            current_chunk.append(word)
            current_size += 1
            
            if current_size >= chunk_size:
                chunks.append(" ".join(current_chunk))
                current_chunk = []
                current_size = 0
        
        # Add remaining words
        if current_chunk:
            chunks.append(" ".join(current_chunk))
        
        logger.info(f"Split text into {len(chunks)} chunks")
        return chunks

# Singleton instance
document_processor = DocumentProcessor()

def get_document_processor():
    return document_processor
